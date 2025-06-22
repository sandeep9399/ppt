# Final Apollo PPT Enhancer (Corrected)

import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io
import requests

st.set_page_config(page_title="Apollo PPT Enhancer", layout="wide")
st.title("🚀 Apollo Slide Enhancer")

st.markdown("""
Upload your **old PowerPoint (.pptx)** file. You'll see AI-powered layout and design suggestions first.
Then, click the button to generate a professionally enhanced Apollo-branded presentation.
""")

uploaded_file = st.file_uploader("Upload old-format PPT", type=["pptx"])
apollo_logo_url = "https://upload.wikimedia.org/wikipedia/en/1/1e/Apollo_Hospitals_Logo.png"

def suggest_design_elements(text):
    text = text.lower()
    if "who" in text:
        return {
            "Layout": "Two-column layout (left: text, right: image)",
            "Font": "Poppins for title, Segoe UI for body", 
            "Color": "Blue-grey with white content boxes",
            "Visual": "WHO quote in callout + real-world image of international healthcare workers",
            "Prompt": "Design a slide with a WHO quote on the left in a callout box. On the right, show a photo of diverse doctors in PPE or lab coats, discussing or reviewing data. Use a blue-grey color palette with Apollo University branding."
        }
    elif "components" in text:
        return {
            "Layout": "4-quadrant grid layout",
            "Font": "Arial Rounded for titles, Segoe UI for text", 
            "Color": "Distinct colors for each quadrant (blue, green, orange, purple)",
            "Visual": "Flat-style infographic showing four labeled sections: Physical, Mental, Social, Spiritual Health with icons",
            "Prompt": "Create an infographic slide divided into 4 color-coded sections labeled Physical, Mental, Social, and Spiritual. Use icons: dumbbell, brain, people, lotus flower respectively. Apollo theme."
        }
    elif "india" in text:
        return {
            "Layout": "India map overlay with side panel",
            "Font": "Segoe UI for clean data visibility",
            "Color": "Warm tones with highlighted states",
            "Visual": "Map of India with health indicators (dots, bars) + key statistics in callouts",
            "Prompt": "Design an infographic map of India showing health data per state. Use hotspot markers, regional shading, and stats callouts in Apollo's visual theme."
        }
    else:
        return {
            "Layout": "Standard title-content layout with visual support",
            "Font": "Segoe UI with Apollo heading styling",
            "Color": "Apollo theme (light backgrounds, blue headers)",
            "Visual": "High-resolution contextual image (e.g., hospital consultation, digital health dashboard)",
            "Prompt": "Create a modern healthcare scene image with medical professionals in discussion using digital tools (tablets, screens) inside a hospital setting, with Apollo-friendly layout."
        }

if uploaded_file:
    old_ppt = Presentation(uploaded_file)
    preview_data = []

    for i, slide in enumerate(old_ppt.slides, start=1):
        text = " ".join([shape.text.strip() for shape in slide.shapes if shape.has_text_frame and shape.text.strip()])
        suggestion = suggest_design_elements(text)
        prompt = f"Design a slide using layout: {suggestion['Layout']}, color: {suggestion['Color']}, with {suggestion['Visual']}. Use Apollo University branding."
        preview_data.append({
            "Slide Number": i,
            "Content Preview": text[:80] + ("..." if len(text) > 80 else ""),
            "Layout": suggestion["Layout"],
            "Font": suggestion["Font"],
            "Color": suggestion["Color"],
            "Visual": suggestion["Visual"],
            "Prompt": prompt
        })

    df = pd.DataFrame(preview_data)
    st.markdown("## Step 1: AI Suggestions Table")
    st.dataframe(df, use_container_width=True)

    st.markdown("## 🔍 Detailed Slide-by-Slide Suggestions")
    for row in preview_data:
        st.markdown(f"**Slide {row['Slide Number']}**")
        st.markdown(f"• **Content Preview**: {row['Content Preview']}")
        st.markdown(f"• **Layout**: {row['Layout']}")
        st.markdown(f"• **Font**: {row['Font']}")
        st.markdown(f"• **Color Theme**: {row['Color']}")
        st.markdown(f"• **Visual Suggestion**: {row['Visual']}")
        st.markdown(f"• **Prompt**: {row['Prompt']}")
        st.markdown("---")

    if st.button("✨ Step 2: Generate Enhanced Apollo Slides"):
        new_ppt = Presentation(uploaded_file)
        new_ppt.slide_width = Inches(13.33)
        new_ppt.slide_height = Inches(7.5)

        title_slide = new_ppt.slides.add_slide(new_ppt.slide_layouts[0])
        title_slide.shapes.title.text = "Apollo University"
        subtitle = title_slide.placeholders[1]
        subtitle.text = "Enhanced Slide Deck — AI Formatted"
        subtitle.text_frame.paragraphs[0].font.name = "Segoe UI"
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        subtitle.text_frame.paragraphs[0].font.italic = True

        # use original slide layout to preserve theme
        layout = old_slide.slide_layout

        for row in preview_data:
            i = row["Slide Number"]
            old_slide = old_ppt.slides[i - 1]
            content_text = [shape.text.strip() for shape in old_slide.shapes if shape.has_text_frame and shape.text.strip()]
            split_required = len(content_text) > 5
            parts = [content_text[:len(content_text)//2], content_text[len(content_text)//2:]] if split_required else [content_text]

            new_slide = new_ppt.slides.add_slide(layout)
            try:
                content_box = new_slide.placeholders[1].text_frame
                content_box.clear()
            except:
                content_box = new_slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(7.5), Inches(4)).text_frame

            if row['Layout'].startswith("4-quadrant"):
                content_box.text = '''🧠 Physical | ❤️ Mental
👥 Social  | 🧘‍♂️ Spiritual

Use a 2x2 layout with icons:
- 🧠 for Physical (dumbbell)
- ❤️ for Mental (brain)
- 👥 for Social (community/people)
- 🧘‍♂️ for Spiritual (lotus or meditation icon)
Use Apollo theme colors in quadrant blocks.'''
            elif row['Layout'].startswith("Two-column"):
                content_box.text = '''🗒️ Left: WHO Quote
"Health is a state of complete physical, mental..." — WHO
📸 Right: Add image of international healthcare team in clinical discussion or data review, wearing lab coats or PPE.
Request image using prompt in speaker notes.'''
            elif row['Layout'].startswith("India map"):
                content_box.text = '''🗺️ Insert map of India
- Overlay hotspots for NCD rates per region
- Use icons for diabetes, cancer, heart, etc.
- Include callouts with % stats
📝 Prompt for AI Image: India health map with disease icons and shaded regions in Apollo colors'''
            else:
                for line in parts[0]:
                    para = content_box.add_paragraph()
                    para.text = line
                    para.font.size = Pt(18)
                    para.font.name = "Segoe UI"
                    para.font.color.rgb = RGBColor(0, 0, 0)


            try:
                content_box = new_slide.placeholders[1].text_frame
            content_box.clear()
        except:
                content_box = new_slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(7.5), Inches(4)).text_frame
            if row['Layout'].startswith("4-quadrant"):
                content_box.text = '''🧠 Physical | ❤️ Mental
👥 Social  | 🧘‍♂️ Spiritual

Use a 2x2 layout with icons:
- 🧠 for Physical (dumbbell)
- ❤️ for Mental (brain)
- 👥 for Social (community/people)
- 🧘‍♂️ for Spiritual (lotus or meditation icon)
Use Apollo theme colors in quadrant blocks.'''
            elif row['Layout'].startswith("Two-column"):
                content_box.text = '''🗒️ Left: WHO Quote
"Health is a state of complete physical, mental..." — WHO
📸 Right: Add image of international healthcare team in clinical discussion or data review, wearing lab coats or PPE.
Request image using prompt in speaker notes.'''
            elif row['Layout'].startswith("India map"):
                content_box.text = '''🗺️ Insert map of India
- Overlay hotspots for NCD rates per region
- Use icons for diabetes, cancer, heart, etc.
- Include callouts with % stats
📝 Prompt for AI Image: India health map with disease icons and shaded regions in Apollo colors'''
            else:
                for line in parts[0]:
                para = content_box.add_paragraph()
                para.text = line
                para.font.size = Pt(18)
                para.font.name = "Segoe UI"
                para.font.color.rgb = RGBColor(0, 0, 0)

            

            
             (note: python-pptx doesn't support native animations)
            animation_note = new_slide.notes_slide.notes_text_frame.add_paragraph()
            animation_note.text = "Suggested Animation: Appear (on click) for each bullet point or image"
            notes_slide = new_slide.notes_slide
            notes_slide.notes_text_frame.text = f'''Full Slide Suggestion:
Layout: {row['Layout']}
Font: {row['Font']}
Color Theme: {row['Color']}
Visual: {row['Visual']}
Prompt: {row['Prompt']}'''

            

            footer = new_slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
            footer_tf = footer.text_frame
            footer_tf.text = "Powered by Apollo Knowledge"
            footer_para = footer_tf.paragraphs[0]
            footer_para.font.name = "Segoe UI"
            footer_para.font.size = Pt(10)
            footer_para.font.italic = True
            footer_para.font.color.rgb = RGBColor(100, 100, 100)

            # image placeholder only; OpenAI call removed
            placeholder_box = new_slide.shapes.add_textbox(Inches(6.0), Inches(1.8), Inches(3.0), Inches(1.5))
            ph_tf = placeholder_box.text_frame
            ph_tf.text = f'''[AI image placeholder]
Prompt:
{row['Prompt']}'''

            if split_required:
                slide2 = new_ppt.slides.add_slide(layout)
                slide2.shapes.title.text = new_slide.shapes.title.text + " (contd)"
                indicator_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.3))
                indicator_tf = indicator_box.text_frame
                indicator_tf.text = "🔁 Continued from previous slide"
                indicator_tf.paragraphs[0].font.size = Pt(14)
                indicator_tf.paragraphs[0].font.color.rgb = RGBColor(90, 90, 90)
                part2_box = slide2.placeholders[1].text_frame
                part2_box.clear()
                for line in parts[1]:
                    para = part2_box.add_paragraph()
                    para.text = line
                    para.font.size = Pt(18)
                    para.font.name = "Segoe UI"
                    para.font.color.rgb = RGBColor(0, 0, 0)
                slide2.notes_slide.notes_text_frame.text = row['Prompt']

        ppt_io = io.BytesIO()
        new_ppt.save(ppt_io)
        ppt_io.seek(0)

        st.download_button(
            label="📥 Download Enhanced Apollo PPT",
            data=ppt_io,
            file_name="Apollo_AI_Enhanced_Slides.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
